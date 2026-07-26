#!/usr/bin/env python3
"""
Capstone Generator V2 — EN/FR templates, Groq AI, 100+ CV layouts
Usage: python generate.py "Name" "email" "Country" "fr|en"
"""
import os, sys, json, random
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from docx import Document
    from docx.shared import Pt as DPt, RGBColor as DRC
    from lxml import etree
except ImportError:
    print("Run: pip install python-pptx python-docx lxml"); sys.exit(1)

from groq_client import generate_content
from cv_engine import build_cv

BASE = Path(__file__).parent
TDIR = BASE.parent / "templates" if (BASE.parent / "templates").exists() else BASE / "templates"
OUT = BASE / "output"
def fill_tf(tf, items):
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.clear()
    for i, (bold, text, bullet_char) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(13)
        r.font.bold = bold
        r.font.name = "Arial"
        
        # Override PPT bullet settings on the paragraph properties level
        pPr = p._p.get_or_add_pPr()
        for tag in ["buChar", "buNone", "buAutoNum", "buBlip"]:
            el = pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag)
            if el is not None:
                pPr.remove(el)
        
        if bullet_char:
            buChar = etree.SubElement(pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar")
            buChar.set("char", bullet_char)
            # Indent and hang for neat bullet alignment
            pPr.set("marL", "360000")
            pPr.set("indent", "-180000")
        else:
            etree.SubElement(pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone")
            pPr.set("marL", "0")
            pPr.set("indent", "0")

# ═══ 1. PRESENTATION ═══
def gen_pptx(name, d, lang):
    template = TDIR / "presentation_template_en.pptx"
    prs = Presentation(str(template))
    o1, o2 = d["opportunity1"], d["opportunity2"]
    
    # 1. Randomize per-slide parameters for rich visual diversity across slides
    header_style = random.choice([1, 2, 3, 4])
    skills_split = random.choice([1, 2, 3])
    
    b_opp = random.choice(["✔", "➤", "📌", "✨"])
    b_skills = random.choice(["✦", "•", "▪", "🔹"])
    b_exp = random.choice(["▸", "❖", "➔", "🔸"])
    b_steps_symbol = random.choice(["➜", "▶", "•", "1."])

    # Per-slide layout diversity (Layout 1: 2-Col, Layout 2: Top-Bottom Stack, Layout 3: Full Width Hero, Layout 4: Left Focus, Layout 5: Right Focus)
    opps_layout = random.choice([1, 2, 3, 4, 5])
    skills_layout = random.choice([1, 2, 3, 4, 5])
    exp_layout = random.choice([1, 2, 3, 4, 5])
    steps_layout = random.choice([1, 2, 3, 4, 5])
    
    # Slide 0: Title — replace Name and Date
    for sh in prs.slides[0].shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    t = r.text.strip()
                    if t == "Name": r.text = name
                    elif t == "Date": r.text = "July 2026"
    
    db = RGBColor(0x1B, 0x3C, 0x6D)

    # Default grid layout coordinates (EMUs)
    top_y = 1152475
    left_x = 311700
    col_w = 3999900
    full_w = 8520600
    right_x = 4832400
    height_y = 3416400

    # Slide 1: Executive Summary & Overview (Dynamic title, 16-17pt readable font, multi-style formatting)
    exec_titles = [
        "Executive Project Summary",
        "Capstone Abstract & Vision",
        "Executive Overview & Strategy",
        "Green Pathways Portfolio Summary",
        "Executive Summary & Placement Objectives"
    ]
    t_exec = random.choice(exec_titles)

    if len(prs.slides) >= 2:
        sl1 = prs.slides[1]
        swtf1 = [(s, s.text_frame) for s in sl1.shapes if s.has_text_frame]
        swtf1.sort(key=lambda x: (x[0].top, x[0].left))
        
        if len(swtf1) >= 2:
            sh_title, tf_title = swtf1[0]
            sh_body, tf_body = swtf1[1]
            
            # Align shape left and top coordinates cleanly to match exact layout grid
            sh_title.left = left_x
            sh_title.top = top_y
            sh_title.width = full_w
            
            sh_body.left = left_x
            sh_body.top = top_y + 1100000
            sh_body.width = full_w
            sh_body.height = height_y - 1100000

            tf_title.margin_left = 0; tf_title.margin_right = 0; tf_title.margin_top = 0; tf_title.margin_bottom = 0
            tf_body.margin_left = 0; tf_body.margin_right = 0; tf_body.margin_top = 0; tf_body.margin_bottom = 0
            
            # 1. Update Title of Slide 1
            tf_title.clear()
            r_title = tf_title.paragraphs[0].add_run()
            r_title.text = t_exec
            r_title.font.size = Pt(20)
            r_title.font.bold = True
            r_title.font.name = "Arial"
            r_title.font.color.rgb = db
            
            pPr = tf_title.paragraphs[0]._p.get_or_add_pPr()
            for tag in ["buChar", "buNone", "buAutoNum", "buBlip"]:
                el = pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag)
                if el is not None:
                    pPr.remove(el)
            etree.SubElement(pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone")

        # 2. Update Body Text of Slide 1 with large font (Pt(16)-Pt(17)) and multi-style formatting
        if len(swtf1) >= 2:
            tf_body.clear()
            tf_body.word_wrap = True
            
            summary_text = d.get("capstone_summary", "")
            sentences = [s.strip() for s in summary_text.split(".") if s.strip()]
            
            s1_format_style = random.choice(["lead_quote", "callout_pillars", "full_paragraph"])
            
            if s1_format_style == "lead_quote" and len(sentences) >= 2:
                # Lead Sentence (Bold & Highlighted, Pt(17))
                p1 = tf_body.paragraphs[0]
                p1.space_after = Pt(10)
                r1 = p1.add_run()
                r1.text = sentences[0] + "."
                r1.font.size = Pt(17)
                r1.font.bold = True
                r1.font.name = "Arial"
                r1.font.color.rgb = db
                
                # Remaining Text (Pt(16))
                p2 = tf_body.add_paragraph()
                p2.space_after = Pt(8)
                r2 = p2.add_run()
                r2.text = " ".join(sentences[1:]) + "."
                r2.font.size = Pt(16)
                r2.font.name = "Arial"
                
            elif s1_format_style == "callout_pillars":
                p1 = tf_body.paragraphs[0]
                p1.space_after = Pt(10)
                r1 = p1.add_run()
                r1.text = summary_text
                r1.font.size = Pt(16)
                r1.font.name = "Arial"
                
                # Pillar Callout line
                p2 = tf_body.add_paragraph()
                p2.space_after = Pt(6)
                r2 = p2.add_run()
                r2.text = f"Pillars of Impact: {d['education']['specialization']} • {d['opportunity1']['company']} • {d['opportunity2']['company']}"
                r2.font.size = Pt(14)
                r2.font.bold = True
                r2.font.name = "Arial"
                r2.font.color.rgb = db
                
            else:
                p1 = tf_body.paragraphs[0]
                p1.space_after = Pt(10)
                r1 = p1.add_run()
                r1.text = summary_text
                r1.font.size = Pt(17)
                r1.font.name = "Arial"

            # Remove bullets on body text paragraphs
            for p in tf_body.paragraphs:
                pPr = p._p.get_or_add_pPr()
                for tag in ["buChar", "buNone", "buAutoNum", "buBlip"]:
                    el = pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag)
                    if el is not None:
                        pPr.remove(el)
                etree.SubElement(pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone")
                pPr.set("marL", "0")
                pPr.set("indent", "0")
    
    db = RGBColor(0x1B, 0x3C, 0x6D)
    opp = "Opportunity"
    pos = "Position"
    loc = "Location"
    dur = "Duration: 6 months (from September 2026)"
    
    # Human-style slide titles
    deck_archetype = random.choice([
        "opportunity_comparison",
        "capstone_technical_deepdive",
        "skills_portfolio_matrix",
        "sector_market_strategy",
        "feasibility_roadmap"
    ])
    
    if deck_archetype == "capstone_technical_deepdive":
        t_opp = random.choice(["Capstone Technical Scope & Target Roles", "Research Focus & Industry Applications", "Technical Problem Statement & Placement Scope"])
        t_skills = random.choice(["Engineering Tool Stack & Analytical Modeling", "Technical Competencies & Simulation Tools", "Energy Systems Toolchain & Skills"])
        t_exp = random.choice(["Capstone Outcomes & Applied Qualifications", "Practical Project Preparation & Assets", "Research Methodology & Applied Experience"])
        t_steps = random.choice(["Technical Implementation Roadmap", "3-Phase Development Strategy", "Next Steps in Capstone Deployment"])
    elif deck_archetype == "skills_portfolio_matrix":
        t_opp = random.choice(["Career Specialization & Target Roles", "Target Internship Placement Overview", "Professional Placement Targets"])
        t_skills = random.choice(["Technical Competencies & Environmental Compliance", "Hard Skills Matrix & Qualifications", "Core Engineering & Applied Skills"])
        t_exp = random.choice(["Practical Project Accomplishments", "Fieldwork & Leadership Experience", "Demonstrated Technical Preparation"])
        t_steps = random.choice(["Professional Growth & Skill Mastery Goals", "Career Strategy & Development Milestones", "Personal Action Plan & Future Goals"])
    elif deck_archetype == "sector_market_strategy":
        t_opp = random.choice(["Regional Clean Energy Placement Landscape", "Sector Analysis & Target Companies", "Clean Tech Market Opportunities"])
        t_skills = random.choice(["Target Role Prerequisites & Core Skills", "Industry Qualifications & Skill Matrix", "Required Competencies & Background"])
        t_exp = random.choice(["Academic Specialization & Applied Competencies", "Relevant Background & Project Portfolio", "Applied Academic Assets"])
        t_steps = random.choice(["Career Deployment & Strategic Milestones", "Strategic Roadmap & Next Steps", "Action Plan for Sector Integration"])
    elif deck_archetype == "feasibility_roadmap":
        t_opp = random.choice(["Strategic Target Organizations & Roles", "Target Placement Scope Overview", "Selected Renewable Energy Opportunities"])
        t_skills = random.choice(["Feasibility Modeling & Energy Analytics", "Technical Modeling & Energy Yield Skills", "Analytical & Feasibility Competencies"])
        t_exp = random.choice(["Value Delivered to Target Employers", "Demonstrated Field & Lab Capabilities", "Qualifications & Practical Assets"])
        t_steps = random.choice(["3-Phase Professional Milestone Strategy", "Strategic Action Roadmap", "Implementation Milestones & Action Plan"])
    else:
        t_opp = random.choice(["Target Internship Opportunities", "Selected Career Pathways", "My Target Roles in Renewable Energy"])
        t_skills = random.choice(["Skills & Core Qualifications", "My Core Competencies", "Technical Skills & Requirements"])
        t_exp = random.choice(["My Relevant Experience", "Practical Background & Transferable Skills", "Projects & Accomplishments"])
        t_steps = random.choice(["My Action Plan & Next Steps", "Career Roadmap & Immediate Actions", "Future Development & Strategy"])
    
    # 2. Format Opportunities text as full narrative paragraphs (bullet_char=None for rich text body)
    if header_style == 1:
        o1_hdr = [(True, f"{opp} 1 : {o1['company']}", None), (False, f"{pos} : {o1['position']}", None), (False, f"{loc} : {o1['location']}", None)]
        o2_hdr = [(True, f"{opp} 2 : {o2['company']}", None), (False, f"{pos} : {o2['position']}", None), (False, f"{loc} : {o2['location']}", None)]
    elif header_style == 2:
        o1_hdr = [(True, f"{o1['position']} @ {o1['company']}", None), (False, f"{loc}: {o1['location']}", None)]
        o2_hdr = [(True, f"{o2['position']} @ {o2['company']}", None), (False, f"{loc}: {o2['location']}", None)]
    elif header_style == 3:
        o1_hdr = [(True, f"{o1['company']} — {o1['position']}", None), (False, f"📍 {o1['location']}", None)]
        o2_hdr = [(True, f"{o2['company']} — {o2['position']}", None), (False, f"📍 {o2['location']}", None)]
    else:
        o1_hdr = [(True, o1['company'], None), (False, f"Role: {o1['position']} ({o1['location']})", None)]
        o2_hdr = [(True, o2['company'], None), (False, f"Role: {o2['position']} ({o2['location']})", None)]
        
    opp_text_style = random.choice(["full_narrative", "callout_block", "clean_paragraph"])
    
    if opp_text_style == "full_narrative":
        o1_desc_text = f"{o1['description']} This target placement aligns directly with my specialized coursework in {d['education']['specialization'].lower()}."
        o2_desc_text = f"{o2['description']} This role provides key practical development in {d['education']['specialization'].lower()}."
    elif opp_text_style == "callout_block":
        o1_desc_text = f"{o1['description']}\nCore Technical Focus: {', '.join(o1['skills'][:2])}."
        o2_desc_text = f"{o2['description']}\nCore Technical Focus: {', '.join(o2['skills'][:2])}."
    else:
        o1_desc_text = o1['description']
        o2_desc_text = o2['description']

    # Set bullet_char=None so description renders as full narrative paragraph body text (not a bulleted item)
    o1_opp_content = o1_hdr + [(False, dur, None), (False, o1_desc_text, None)]
    o2_opp_content = o2_hdr + [(False, dur, None), (False, o2_desc_text, None)]

    # 3. Format Skills & Qualifications based on skills_split style (bold, text, bullet_char)
    q_label = random.choice(["Academic Requirements", "Required Qualifications", "Minimum Prerequisites", "Target Education Level"])
    
    if skills_split == 1:
        lbl_s1 = random.choice([f"Core Skills for {o1['company']}", f"Technical Competencies - {o1['company']}", f"Skills Required by {o1['company']}"])
        lbl_s2 = random.choice([f"Core Skills for {o2['company']}", f"Technical Competencies - {o2['company']}", f"Skills Required by {o2['company']}"])
        skills_left = [(True, lbl_s1, None)] + [(False, s, b_skills) for s in o1['skills']] + [(True, q_label, None), (False, o1['qualification'], b_skills)]
        skills_right = [(True, lbl_s2, None)] + [(False, s, b_skills) for s in o2['skills']] + [(True, q_label, None), (False, o2['qualification'], b_skills)]
    elif skills_split == 2:
        tech_skills = sorted(list(set(o1['skills'] + o2['skills'])))[:5]
        soft_skills = d.get("cv_skills", {}).get("right", ["Teamwork", "Communication", "Problem-solving"])
        lbl_tech = random.choice(["Technical Skills Checklist", "My Technical Skills", "Solar & Energy Skills", "Target Hard Skills"])
        lbl_soft = random.choice(["Professional Soft Skills", "Core Competencies", "Key Transferable Skills", "Interpersonal Strengths"])
        skills_left = [(True, lbl_tech, None)] + [(False, s, b_skills) for s in tech_skills] + [(True, "Education Requirements", None), (False, o1['qualification'], b_skills), (False, o2['qualification'], b_skills)]
        skills_right = [(True, lbl_soft, None)] + [(False, s, b_skills) for s in soft_skills]
    else:
        all_skills = sorted(list(set(o1['skills'] + o2['skills'])))[:6]
        lbl_all = random.choice(["Combined Skills Checklist", "Prerequisites & Core Competencies", "Full Skills Profile"])
        lbl_roles = random.choice(["Role Alignment Details", "Target Internships & Positions", "Internship Opportunities"])
        skills_left = [(True, lbl_all, None)] + [(False, s, b_skills) for s in all_skills]
        skills_right = [(True, lbl_roles, None), (True, o1['company'], None), (False, o1['position'], b_skills), (True, o2['company'], None), (False, o2['position'], b_skills)]

    # Content slide definitions with individual layout types
    lbl_exp1 = random.choice([f"Preparation for {o1['company']}", f"Skill Alignment for {o1['company']}", f"Target Assets for {o1['company']}", f"Qualifications for {o1['company']}"])
    lbl_exp2 = random.choice([f"Preparation for {o2['company']}", f"Skill Alignment for {o2['company']}", f"Target Assets for {o2['company']}", f"Qualifications for {o2['company']}"])
    
    exp_slide = {
        "t": t_exp,
        "l": [(True, lbl_exp1, None)] + [(False, e, b_exp) for e in o1['my_experience']],
        "r": [(True, lbl_exp2, None)] + [(False, e, b_exp) for e in o2['my_experience']],
        "layout": exp_layout
    }
    
    lbl_step1 = random.choice(["Immediate Actions (Summer 2026)", "Short-Term Action Items", "Priority Action Plan", "Immediate Next Steps"])
    lbl_step2 = random.choice(["Medium-Term Development", "Professional Development Goals", "Long-Term Objectives", "Development Strategy"])
    
    def fmt_steps(steps, b_char):
        if b_char == "1.":
            return [(False, f"{idx+1}. {x}", None) for idx, x in enumerate(steps)]
        return [(False, x, b_char) for x in steps]

    steps_slide = {
        "t": t_steps,
        "l": [(True, lbl_step1, None)] + fmt_steps(d['next_steps_immediate'], b_steps_symbol),
        "r": [(True, lbl_step2, None)] + fmt_steps(d['next_steps_medium'], b_steps_symbol),
        "layout": steps_layout
    }

    opps_slide = {"t": t_opp, "l": o1_opp_content, "r": o2_opp_content, "layout": opps_layout}
    skills_slide = {"t": t_skills, "l": skills_left, "r": skills_right, "layout": skills_layout}

    slide_content = [opps_slide, skills_slide, exp_slide, steps_slide]
    
    # 4. Apply Column Order Swapping dynamically per slide (50% chance)
    for s_dict in slide_content:
        if random.choice([True, False]):
            s_dict["l"], s_dict["r"] = s_dict["r"], s_dict["l"]
            
    ci = [4, 6, 8, 10]
    layout_sequence = []
    
    # Default coordinates (EMUs)
    top_y = 1152475
    left_x = 311700
    col_w = 3999900
    full_w = 8520600
    right_x = 4832400
    height_y = 3416400
    
    for idx, c in zip(ci, slide_content):
        if idx >= len(prs.slides): continue
        sl = prs.slides[idx]
        swtf = [(s, s.text_frame) for s in sl.shapes if s.has_text_frame]
        swtf.sort(key=lambda x: (x[0].top, x[0].left))
        
        if len(swtf) >= 3:
            # Title positioning and text replacement
            tf = swtf[0][1]; tf.clear()
            r = tf.paragraphs[0].add_run(); r.text = c["t"]
            r.font.size = Pt(20); r.font.bold = True; r.font.name = "Arial"; r.font.color.rgb = db
            
            # Reset bullet settings for slide titles
            pPr = tf.paragraphs[0]._p.get_or_add_pPr()
            for tag in ["buChar", "buNone", "buAutoNum", "buBlip"]:
                el = pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag)
                if el is not None:
                    pPr.remove(el)
            etree.SubElement(pPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone")
            pPr.set("marL", "0")
            pPr.set("indent", "0")

            bs = sorted(swtf[1:], key=lambda x: x[0].left)
            if len(bs) >= 2:
                s1, s2 = bs[0][0], bs[1][0]
                cur_layout = c["layout"]
                
                # Record the layout sequence for this specific slide
                layout_sequence.append(f"L{cur_layout}")
                
                # Delete old placeholder shapes to prevent master slide styling overrides
                sl.shapes._element.remove(s1._element)
                sl.shapes._element.remove(s2._element)
                
                if cur_layout == 1:
                    # Side-by-side Columns
                    tb1 = sl.shapes.add_textbox(left_x, top_y, col_w, height_y)
                    tb2 = sl.shapes.add_textbox(right_x, top_y, col_w, height_y)
                    fill_tf(tb1.text_frame, c["l"])
                    fill_tf(tb2.text_frame, c["r"])
                    
                elif cur_layout == 2:
                    # Top-and-Bottom Split (Vertical Stack)
                    tb1 = sl.shapes.add_textbox(left_x, top_y, full_w, 1600000)
                    tb2 = sl.shapes.add_textbox(left_x, 2900000, full_w, 1600000)
                    fill_tf(tb1.text_frame, c["l"])
                    fill_tf(tb2.text_frame, c["r"])
                    
                elif cur_layout == 3:
                    # Single Column Full-Width
                    tb1 = sl.shapes.add_textbox(left_x, top_y, full_w, height_y)
                    # Merge content from both columns
                    merged = c["l"] + [(True, "", None)] + c["r"]
                    fill_tf(tb1.text_frame, merged)
                    
                elif cur_layout == 4:
                    # Asymmetric (Left Highlight / Right List)
                    tb1 = sl.shapes.add_textbox(left_x, top_y, 2800000, height_y)
                    tb2 = sl.shapes.add_textbox(3400000, top_y, 5432000, height_y)
                    fill_tf(tb1.text_frame, c["l"])
                    fill_tf(tb2.text_frame, c["r"])
                    
                else:
                    # Asymmetric (Left List / Right Highlight)
                    tb1 = sl.shapes.add_textbox(left_x, top_y, 5432000, height_y)
                    tb2 = sl.shapes.add_textbox(6032300, top_y, 2800000, height_y)
                    fill_tf(tb1.text_frame, c["l"])
                    fill_tf(tb2.text_frame, c["r"])
    
    # Clean metadata - set author to user
    prs.core_properties.author = name
    prs.core_properties.last_modified_by = name
    prs.core_properties.title = f'Green Pathways Capstone - {name}'
    prs.core_properties.comments = ''
    prs.core_properties.subject = 'Green Pathways Capstone Project'
    
    # Safe PPT design ID for filename compatibility (no special chars like arrow, check, diamond in filenames)
    bullet_names = {"•": "dot", "✔": "check", "✦": "star", "➤": "arrow", "❖": "diamond"}
    bullet_name = bullet_names.get(b_opp, "bullet")
    layout_code = "".join(layout_sequence)
    ppt_design_safe = f"H{header_style}_B_{bullet_name}_S{skills_split}_{layout_code}"
    
    ppt_patterns = [
        "NAME_Capstone_Presentation.pptx",
        "NAME_GreenPathways_Portfolio.pptx",
        "Capstone_Project_Presentation_NAME.pptx",
        "NAME_Renewable_Energy_Slides.pptx",
        "NAME_Project_Pitch.pptx",
        "Green_Pathways_Capstone_NAME.pptx"
    ]
    safe_name = name.replace(' ', '_')
    fn = random.choice(ppt_patterns).replace("NAME", safe_name)
    out = OUT / fn
    prs.save(str(out)); return out, ppt_design_safe

# ═══ 2. CAPSTONE DOCUMENT ═══
def gen_docx(name, d, lang):
    template = TDIR / "capstone_template_en.docx"
    doc = Document(str(template))
    
    doc_fonts = ["Sora", "Calibri", "Arial", "Georgia", "Trebuchet MS", "Century Gothic", "Palatino Linotype"]
    doc_colors = [
        DRC(0x15, 0x39, 0x88), # Classic Blue
        DRC(0x1A, 0x47, 0x2A), # Pine Green
        DRC(0x2C, 0x3E, 0x50), # Slate
        DRC(0x3C, 0x15, 0x18), # Deep Wine
        DRC(0x2D, 0x34, 0x36), # Graphite
        DRC(0x0C, 0x35, 0x47), # Deep Sea
        DRC(0x1B, 0x43, 0x32)  # Emerald
    ]
    fn_name = random.choice(doc_fonts)
    fn_color = random.choice(doc_colors)
    
    # Put name on RIGHT side of table
    for table in doc.tables:
        for row in table.rows:
            cells = row.cells
            if len(cells) >= 2:
                left = cells[0].text.strip()
                if "Your Name" in left or "Name" in left:
                    cells[1].paragraphs[0].clear()
                    r = cells[1].paragraphs[0].add_run(name)
                    r.font.name = fn_name; r.font.size = DPt(14)
                    r.italic = True; r.font.color.rgb = fn_color
                    break
    
    # Fill Table 1 = summary, Table 3 = CV improvements
    tables = doc.tables
    if len(tables) >= 4:
        cell = tables[1].rows[0].cells[0]
        cell.paragraphs[0].clear()
        r = cell.paragraphs[0].add_run(d["capstone_summary"])
        r.font.name = fn_name; r.font.size = DPt(10.5)
        
        cell = tables[3].rows[0].cells[0]
        cell.paragraphs[0].clear()
        r = cell.paragraphs[0].add_run(d["cv_improvements"])
        r.font.name = fn_name; r.font.size = DPt(10.5)
    
    # Clean metadata
    doc.core_properties.author = name
    doc.core_properties.last_modified_by = name
    doc.core_properties.title = f'Green Pathways Capstone - {name}'
    doc.core_properties.comments = ''
    
    doc_patterns = [
        "NAME_Capstone_Document.docx",
        "NAME_Capstone_Summary_Report.docx",
        "Capstone_Summary_NAME.docx",
        "NAME_GreenPathways_Report.docx",
        "Project_Summary_Report_NAME.docx",
        "NAME_Portfolio_Document.docx"
    ]
    safe_name = name.replace(' ', '_')
    fn = random.choice(doc_patterns).replace("NAME", safe_name)
    out = OUT / fn
    doc.save(str(out)); return out

# ═══ MAIN ═══
def main():
    if len(sys.argv) >= 4:
        name, email, country = sys.argv[1], sys.argv[2], sys.argv[3]
    elif len(sys.argv) >= 3:
        name, email, country = sys.argv[1], sys.argv[2], "France"
    else:
        print('Usage: python generate.py "Name" "email" "Country"'); sys.exit(1)
    
    name = " ".join(name.split())
    lang = "en"
    
    flags = {"France": "🇫🇷", "Belgique": "🇧🇪", "Ghana": "🇬🇭", "Ireland": "🇮🇪"}
    print(f"\n{'='*50}")
    print(f"🌱 {name} | {flags.get(country, '')} {country} | EN")
    print(f"{'='*50}")
    
    print("\n🌱 Generating content from database...")
    data = generate_content(name, email, country, lang)
    print(f"   ✅ Content ready")
    print(f"   🏢 {data['opportunity1']['company']} + {data['opportunity2']['company']}")
    print(f"   🎓 {data['education']['university']}")
    
    OUT.mkdir(exist_ok=True)
    
    print(f"\n📊 [1/3] Presentation ({lang.upper()})...")
    p1_path, ppt_design = gen_pptx(name, data, lang); print(f"   ✅ {p1_path.name}")
    
    print(f"📄 [2/3] Capstone Document ({lang.upper()})...")
    p2 = gen_docx(name, data, lang); print(f"   ✅ {p2.name}")
    
    print("📝 [3/3] CV...")
    temp_path = OUT / f"CV_{name.replace(' ','_')}_temp.docx"
    design = build_cv(data, temp_path)
    cv_patterns = [
        "NAME_Resume.docx",
        "NAME_CV.docx",
        "Curriculum_Vitae_NAME.docx",
        "NAME_Green_Energy_CV.docx",
        "Resume_NAME.docx",
        "CV_NAME_Renewable_Energy.docx"
    ]
    safe_name = name.replace(' ', '_')
    fn = random.choice(cv_patterns).replace("NAME", safe_name)
    cv_path = OUT / fn
    if temp_path.exists():
        temp_path.rename(cv_path)
    print(f"   🎨 Layout: {design}")
    print(f"   ✅ {cv_path.name}")
    
    # Save content for preview
    preview = {
        "name": name, "email": data.get("email",""), "country": data.get("country",""),
        "phone": data.get("phone",""), "city": data.get("city",""), "lang": lang,
        "opp1": data.get("opportunity1",{}), "opp2": data.get("opportunity2",{}),
        "education": data.get("education",{}), "experience": data.get("experience",{}),
        "profile": data.get("profile_summary",""), "skills": data.get("cv_skills",{}),
        "languages": data.get("languages",[]), "interests": data.get("interests",""),
        "summary": data.get("capstone_summary",""), "cv_improvements": data.get("cv_improvements",""),
        "next_immediate": data.get("next_steps_immediate",[]), "next_medium": data.get("next_steps_medium",[]),
        "assignment_title": data.get("assignment_title",""),
        "assignment_description": data.get("assignment_description",""),
        "cv_design": design,
        "ppt_design": ppt_design,
    }
    with open(OUT / "content.json", "w") as f:
        json.dump(preview, f, ensure_ascii=False, indent=2)
    
    print(f"\n🎉 Done! → {OUT}/\n")

if __name__ == "__main__":
    main()